import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import io

# Interface Streamlit pour charger les fichiers
st.title("Application Excel -> Powerpoint ")
uploaded_file = st.file_uploader("Uploader un fichier Excel", type=["xlsx"])
uploaded_template = st.file_uploader("Uploader le template PowerPoint", type=["pptx"])

# Vérifier que les deux fichiers sont chargés
if uploaded_file is not None and uploaded_template is not None:
    # Lire le fichier Excel
    dff = pd.read_excel(uploaded_file)

    # Sélectionner une entreprise à partir de la colonne 'POEI'
    companies = dff['POEI'].unique()
    selected_company = st.selectbox("Sélectionner une entreprise", companies)

    if selected_company:
        # Filtrer les données pour l'entreprise sélectionnée
        row_values = dff[dff['POEI'] == selected_company].iloc[0]

        # Définir les colonnes pour les graphiques
        pie_columns = ["candidats contactés par AJC", " PROFILS KO", " NRP",
                       "candidats écartés par AJC", "CANDIDATURES ENVOYEES"]
        bar_columns = [" CV TRAITES", "Candidats profil jugé KO par AJC", "candidats contactés par AJC",
                       "candidats écartés par AJC", "CANDIDATURES ENVOYEES"]
        bar_columns_sixieme_slide = ["CANDIDATURES ENVOYEES", "PROFILS KO SUR CANDIDATURES ENVOYEES", 'VALIDES',
                                     'Nbre entrée en formation', 'Nbre présent en formation J1']

        # Extraire et convertir les valeurs pour les graphiques (remplacer NaN par 0)
        pie_values = {col: pd.to_numeric(row_values.get(col, 0), errors='coerce') for col in pie_columns if col in row_values}
        pie_values = {k: v if not pd.isna(v) else 0 for k, v in pie_values.items()}
        pie_labels = list(pie_values.keys())
        pie_data = list(pie_values.values())

        bar_values = {col: pd.to_numeric(row_values.get(col, 0), errors='coerce') for col in bar_columns if col in row_values}
        bar_values = {k: v if not pd.isna(v) else 0 for k, v in bar_values.items()}
        bar_labels = list(bar_values.keys())
        bar_data = list(bar_values.values())

        bar_values_sixieme_slide = {col: pd.to_numeric(row_values.get(col, 0), errors='coerce') for col in bar_columns_sixieme_slide if col in row_values}
        bar_values_sixieme_slide = {k: v if not pd.isna(v) else 0 for k, v in bar_values_sixieme_slide.items()}
        bar_labels_sixieme_slide = list(bar_values_sixieme_slide.keys())
        bar_data_sixieme_slide = list(bar_values_sixieme_slide.values())

        # Extraire les valeurs textuelles (remplacer NaN par 0)
        a_traite = int(row_values[" CV TRAITES"]) if pd.notna(row_values[" CV TRAITES"]) else 0
        a_contacte = int(row_values["candidats contactés par AJC"]) if pd.notna(row_values["candidats contactés par AJC"]) else 0
        a_retenu = int(row_values["CANDIDATURES ENVOYEES"]) if pd.notna(row_values["CANDIDATURES ENVOYEES"]) else 0
        a_ecarte = int(row_values["candidats écartés par AJC"]) if pd.notna(row_values["candidats écartés par AJC"]) else 0
        entreprise_a_ecarte = int(row_values["PROFILS KO SUR CANDIDATURES ENVOYEES"]) if pd.notna(row_values["PROFILS KO SUR CANDIDATURES ENVOYEES"]) else 0
        entreprise_a_retenu = int(row_values["VALIDES"]) if pd.notna(row_values["VALIDES"]) else 0
        nb_entree_formation = int(row_values["Nbre entrée en formation"]) if pd.notna(row_values["Nbre entrée en formation"]) else 0
        nb_entree_formation_j1 = int(row_values["Nbre présent en formation J1"]) if pd.notna(row_values["Nbre présent en formation J1"]) else 0

        # Bouton pour générer le PowerPoint
        if st.button("Mettre à jour et télécharger PowerPoint"):
            # Charger le template PowerPoint
            prs = Presentation(uploaded_template)

            # Mettre à jour les textes dans les diapositives
            slide_text = prs.slides[3]  # Diapositive 4
            text_first = prs.slides[0]  # Diapositive 1
            text_second = prs.slides[1] # Diapositive 2
            text_six = prs.slides[5]    # Diapositive 6
            last_slide = prs.slides[9]


            # Diapositive 4
            for shape in slide_text.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        if "a traité" in paragraph.text:
                            paragraph.text = f"{'a traité :'.ljust(20)}{a_traite} candidatures"  # Alignement avec 20 caractères au total
                        elif "a contacté" in paragraph.text:
                            paragraph.text = f"{'a contacté :'.ljust(20)}{a_contacte} candidats"
                        elif "a retenu" in paragraph.text:
                            paragraph.text = f"{'a retenu :'.ljust(20)}{a_retenu} candidats"
                        elif "a écarté" in paragraph.text:
                            paragraph.text = f"{'a écarté :'.ljust(20)}{a_ecarte} candidats"
                        elif "Sur l’action" in paragraph.text:
                            paragraph.text = f"Sur l’action {selected_company} du MOIS ANNEE"

            # Diapositive 1
            for shape in text_first.shapes:
                if shape.has_text_frame:
                    text_frame_first = shape.text_frame
                    for paragraph in text_frame_first.paragraphs:
                        if "Bilan des actions POEI*" in paragraph.text:
                            paragraph.text = f"Bilan des actions POEI* {selected_company}"

            # Diapositive 2
            for shape in text_second.shapes:
                if shape.has_text_frame:
                    text_frame_second = shape.text_frame
                    for paragraph in text_frame_second.paragraphs:
                        if "nom_entreprise Session mois année" in paragraph.text:
                            paragraph.text = f"{selected_company}"
                        elif "Sur l’action de recrutement de l’action POEI" in paragraph.text:
                            paragraph.text = f"Sur l’action de recrutement de l’action POEI {selected_company}"

            # Diapositive 6
            for shape in text_six.shapes:
                if shape.has_text_frame:
                    text_frame_six = shape.text_frame
                    for paragraph in text_frame_six.paragraphs:
                        if "Candidatures envoyées par AJC à" in paragraph.text:
                            paragraph.text = f"Candidatures envoyées par AJC à {selected_company}"
                        elif "Sur l’action de recrutement de l’action POEI " in paragraph.text:
                            paragraph.text = f"Sur l’action de recrutement de l’action POEI {selected_company}"
                        elif "a envoyé" in paragraph.text:
                            paragraph.text = f"a envoyé {a_retenu} candidatures à {selected_company}"
                        elif "a écarté" in paragraph.text:
                            paragraph.text = f"a écarté {entreprise_a_ecarte} \n (CV KO, Entretien KO, Test KO, Process trop long, désistement)"
                        elif "a retenu" in paragraph.text:
                            paragraph.text = f"a retenu {entreprise_a_retenu}"

            # Ajouter un graphique en camembert dans la diapositive 4 (index 3)
            chart_data_pie = CategoryChartData()
            chart_data_pie.categories = pie_labels
            chart_data_pie.add_series('Série 1', pie_data)
            slide_pie = prs.slides[3]
            x, y, cx, cy = Inches(1), Inches(1), Inches(4.5), Inches(4.5)
            chart_pie = slide_pie.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_pie
            ).chart

            # Ajouter un graphique en barres dans la diapositive 5 (index 4)
            chart_data_bar = CategoryChartData()
            chart_data_bar.categories = bar_labels
            chart_data_bar.add_series('Série 1', bar_data)
            slide_bar = prs.slides[4]
            x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(4.5)
            chart_bar = slide_bar.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_bar
            ).chart

            # Ajouter un second graphique en barres dans la diapositive 6 (index 5)
            chart_data_bar_six = CategoryChartData()
            chart_data_bar_six.categories = bar_labels_sixieme_slide
            chart_data_bar_six.add_series('Série 1', bar_data_sixieme_slide)
            slide_bar_six = prs.slides[5]
            x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(4.5)
            chart_bar_six = slide_bar_six.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_bar_six
            ).chart

            # Sauvegarder le PowerPoint dans un objet BytesIO
            output_io = io.BytesIO()
            prs.save(output_io)
            output_io.seek(0)

            # Bouton de téléchargement
            st.download_button(
                label="Télécharger PowerPoint mis à jour",
                data=output_io,
                file_name=f"presentation_{selected_company}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )