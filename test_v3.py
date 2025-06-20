import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import io

# Interface Streamlit pour charger les fichiers
st.title("Application Excel -> PowerPoint")
uploaded_file = st.file_uploader("Uploader un fichier Excel", type=["xlsx"])
uploaded_template = st.file_uploader("Uploader le template PowerPoint (10 diapositives vides)", type=["pptx"])

# Vérifier que les deux fichiers sont chargés
if uploaded_file is not None and uploaded_template is not None:
    # Lire le fichier Excel
    dff = pd.read_excel(uploaded_file)

    # Sélectionner une entreprise à partir de la colonne 'POEI'
    companies = dff['POEI'].unique()
    selected_company = st.selectbox("Sélectionner une entreprise", companies)

    if selected_company:
        # Filtrer les données pour l'entreprise sélectionnée
        row_values = dff[dff['POEI'] == selected_company].iloc[0]

        # Définir les colonnes pour les graphiques
        pie_columns = ["candidats contactés par AJC", " PROFILS KO", " NRP",
                       "candidats écartés par AJC", "CANDIDATURES ENVOYEES"]
        bar_columns = [" CV TRAITES", "Candidats profil jugé KO par AJC", "candidats contactés par AJC",
                       "candidats écartés par AJC", "CANDIDATURES ENVOYEES"]
        bar_columns_sixieme_slide = ["CANDIDATURES ENVOYEES", "PROFILS KO SUR CANDIDATURES ENVOYEES", 'VALIDES',
                                     'Nbre entrée en formation', 'Nbre présent en formation J1']

        # Extraire et convertir les valeurs pour les graphiques (remplacer NaN par 0)
        pie_values = {col: pd.to_numeric(row_values.get(col, 0), errors='coerce') for col in pie_columns if col in row_values}
        pie_values = {k: v if not pd.isna(v) else 0 for k, v in pie_values.items()}
        pie_labels = list(pie_values.keys())
        pie_data = list(pie_values.values())

        bar_values = {col: pd.to_numeric(row_values.get(col, 0), errors='coerce') for col in bar_columns if col in row_values}
        bar_values = {k: v if not pd.isna(v) else 0 for k, v in bar_values.items()}
        bar_labels = list(bar_values.keys())
        bar_data = list(bar_values.values())

        bar_values_sixieme_slide = {col: pd.to_numeric(row_values.get(col, 0), errors='coerce') for col in bar_columns_sixieme_slide if col in row_values}
        bar_values_sixieme_slide = {k: v if not pd.isna(v) else 0 for k, v in bar_values_sixieme_slide.items()}
        bar_labels_sixieme_slide = list(bar_values_sixieme_slide.keys())
        bar_data_sixieme_slide = list(bar_values_sixieme_slide.values())

        # Extraire les valeurs textuelles (remplacer NaN par 0)
        a_traite = int(row_values[" CV TRAITES"]) if pd.notna(row_values[" CV TRAITES"]) else 0
        a_contacte = int(row_values["candidats contactés par AJC"]) if pd.notna(row_values["candidats contactés par AJC"]) else 0
        a_retenu = int(row_values["CANDIDATURES ENVOYEES"]) if pd.notna(row_values["CANDIDATURES ENVOYEES"]) else 0
        a_ecarte = int(row_values["candidats écartés par AJC"]) if pd.notna(row_values["candidats écartés par AJC"]) else 0
        entreprise_a_ecarte = int(row_values["PROFILS KO SUR CANDIDATURES ENVOYEES"]) if pd.notna(row_values["PROFILS KO SUR CANDIDATURES ENVOYEES"]) else 0
        entreprise_a_retenu = int(row_values["VALIDES"]) if pd.notna(row_values["VALIDES"]) else 0
        nb_entree_formation = int(row_values["Nbre entrée en formation"]) if pd.notna(row_values["Nbre entrée en formation"]) else 0
        nb_entree_formation_j1 = int(row_values["Nbre présent en formation J1"]) if pd.notna(row_values["Nbre présent en formation J1"]) else 0

        # Bouton pour générer le PowerPoint
        if st.button("Mettre à jour et télécharger PowerPoint"):
            # Charger le template PowerPoint (avec 10 diapositives vides)
            prs = Presentation(uploaded_template)

            # Définir les 10 diapositives
            slide_1 = prs.slides[0]  # Diapositive 1
            slide_2 = prs.slides[1]  # Diapositive 2
            slide_3 = prs.slides[2]  # Diapositive 3
            slide_4 = prs.slides[3]  # Diapositive 4
            slide_5 = prs.slides[4]  # Diapositive 5
            slide_6 = prs.slides[5]  # Diapositive 6
            slide_7 = prs.slides[6]  # Diapositive 7
            slide_8 = prs.slides[7]  # Diapositive 8
            slide_9 = prs.slides[8]  # Diapositive 9
            slide_10 = prs.slides[9]  # Diapositive 10

            # --- Diapositive 1 : Titre principal ---
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(1)
            textbox_1 = slide_1.shapes.add_textbox(left, top, width, height)
            text_frame_1 = textbox_1.text_frame
            p = text_frame_1.add_paragraph()
            p.text = f"Bilan des actions POEI* {selected_company}"
            p.font.name = "Poppins"
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # --- Diapositive 2 : Rapport d’activité ---
            textbox_2a = slide_2.shapes.add_textbox(left, top, width, height)
            text_frame_2a = textbox_2a.text_frame
            p = text_frame_2a.add_paragraph()
            p.text = "Rapport d’activité recrutement"
            p.font.name = "Poppins"
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            textbox_2b = slide_2.shapes.add_textbox(left, top + Inches(1.5), width, height)
            text_frame_2b = textbox_2b.text_frame
            p = text_frame_2b.add_paragraph()
            p.text = f"{selected_company}"
            p.font.name = "Poppins"
            p.font.size = Pt(44)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # --- Diapositive 3 : Introduction générale ---
            textbox_3 = slide_3.shapes.add_textbox(left, top, width, height)
            text_frame_3 = textbox_3.text_frame
            p = text_frame_3.add_paragraph()
            p.text = "Introduction aux activités de recrutement d'AJC"
            p.font.name = "Poppins"
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # --- Diapositive 4 : Statistiques et camembert ---
            textbox_4a = slide_4.shapes.add_textbox(left, top, width, height)
            text_frame_4a = textbox_4a.text_frame
            p = text_frame_4a.add_paragraph()
            p.text = "Candidatures Traitées / Ecartées par AJC"
            p.font.name = "Poppins"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 128)  # Bleu marine
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            textbox_4b = slide_4.shapes.add_textbox(left, top + Inches(0.5), width, height)
            text_frame_4b = textbox_4b.text_frame
            p = text_frame_4b.add_paragraph()
            p.text = f"Sur l’action {selected_company} du MOIS ANNEE"
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 0, 0)  # Noir
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            stats_text = (
                f"a traité :          {a_traite} candidatures\n"
                f"a contacté :        {a_contacte} candidats\n"
                f"a retenu :          {a_retenu} candidats\n"
                f"a écarté :          {a_ecarte} candidats\n"
                "Soit un taux de sélection de 13% de l’ensemble des candidatures traitées"
            )
            textbox_4c = slide_4.shapes.add_textbox(left, top + Inches(1.5), width, Inches(3))
            text_frame_4c = textbox_4c.text_frame
            p = text_frame_4c.add_paragraph()
            p.text = stats_text
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # Ajouter le graphique en camembert
            chart_data_pie = CategoryChartData()
            chart_data_pie.categories = pie_labels
            chart_data_pie.add_series('Série 1', pie_data)
            x, y, cx, cy = Inches(5), Inches(1), Inches(4), Inches(4)
            chart_pie = slide_4.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_pie).chart

            # --- Diapositive 5 : Graphique en barres 1 ---
            chart_data_bar = CategoryChartData()
            chart_data_bar.categories = bar_labels
            chart_data_bar.add_series('Série 1', bar_data)
            x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(4.5)
            chart_bar = slide_5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_bar).chart

            # --- Diapositive 6 : Statistiques détaillées et graphique en barres 2 ---
            textbox_6a = slide_6.shapes.add_textbox(left, top, width, height)
            text_frame_6a = textbox_6a.text_frame
            p = text_frame_6a.add_paragraph()
            p.text = f"Candidatures envoyées par AJC à {selected_company}"
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            textbox_6b = slide_6.shapes.add_textbox(left, top + Inches(0.5), width, height)
            text_frame_6b = textbox_6b.text_frame
            p = text_frame_6b.add_paragraph()
            p.text = f"Sur l’action de recrutement de l’action POEI {selected_company}"
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            textbox_6c = slide_6.shapes.add_textbox(left, top + Inches(1.0), width, height)
            text_frame_6c = textbox_6c.text_frame
            p = text_frame_6c.add_paragraph()
            p.text = f"a envoyé {a_retenu} candidatures à {selected_company}"
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            textbox_6d = slide_6.shapes.add_textbox(left, top + Inches(1.5), width, height)
            text_frame_6d = textbox_6d.text_frame
            p = text_frame_6d.add_paragraph()
            p.text = f"a écarté {entreprise_a_ecarte} \n (CV KO, Entretien KO, Test KO, Process trop long, désistement)"
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            textbox_6e = slide_6.shapes.add_textbox(left, top + Inches(2.0), width, height)
            text_frame_6e = textbox_6e.text_frame
            p = text_frame_6e.add_paragraph()
            p.text = f"a retenu {entreprise_a_retenu}"
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # Ajouter le second graphique en barres
            x, y, cx, cy = Inches(5), Inches(1), Inches(4), Inches(4)
            chart_data_bar_six = CategoryChartData()
            chart_data_bar_six.categories = bar_labels_sixieme_slide
            chart_data_bar_six.add_series('Série 1', bar_data_sixieme_slide)
            chart_bar_six = slide_6.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_bar_six).chart

            # --- Diapositive 10 : Conclusion avec contacts ---
            textbox_10 = slide_10.shapes.add_textbox(left, top, width,
                                                     Inches(2))  # Hauteur augmentée pour tout le texte
            text_frame_10 = textbox_10.text_frame
            contact_text = (
                "À votre disposition pour discuter du futur\n"
                "Contact\n"
                "6 rue Rougemont, 75009 Paris\n"
                "Tel : 01 75 43 86 72\n"
                "formonsnous@ajc-ingenierie.fr\n"
                "www.ajc-ingenierie.fr\n"
                "www.unjourunjob.fr"
            )
            p = text_frame_10.add_paragraph()
            p.text = contact_text
            p.font.name = "Poppins"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanc
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # Sauvegarder le PowerPoint dans un objet BytesIO
            output_io = io.BytesIO()
            prs.save(output_io)
            output_io.seek(0)

            # Bouton de téléchargement
            st.download_button(
                label="Télécharger PowerPoint mis à jour",
                data=output_io,
                file_name=f"presentation_{selected_company}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )